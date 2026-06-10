"""
PowerPoint generation service.

Replaces the Gamma API approach with a local pipeline:
master analysis → content extraction → variant generation via Claude (Bedrock) → PPTX assembly.
"""

from __future__ import annotations

import json
import os
import tempfile
import uuid
from pathlib import Path
from typing import Any

import fitz
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# --- LLM Configuration: Bedrock (primary) or Anthropic API (fallback) ---

BEDROCK_AWS_ACCESS_KEY_ID = os.environ.get("BEDROCK_AWS_ACCESS_KEY_ID", "")
BEDROCK_AWS_SECRET_ACCESS_KEY = os.environ.get("BEDROCK_AWS_SECRET_ACCESS_KEY", "")
BEDROCK_AWS_REGION = os.environ.get("BEDROCK_AWS_REGION", "eu-west-1")
BEDROCK_INFERENCE_PREFIX = os.environ.get("BEDROCK_INFERENCE_PREFIX", "eu")
BEDROCK_MODEL = os.environ.get(
    "BEDROCK_MODEL",
    f"{BEDROCK_INFERENCE_PREFIX}.anthropic.claude-opus-4-6-v1"
)

USE_BEDROCK = bool(BEDROCK_AWS_ACCESS_KEY_ID and BEDROCK_AWS_SECRET_ACCESS_KEY)

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
ANTHROPIC_MODEL = "claude-sonnet-4-20250514"
ANTHROPIC_API_URL = "https://api.anthropic.com/v1/messages"
ANTHROPIC_VERSION = "2023-06-01"

_PLACEHOLDER_TYPE_MAP = {
    PP_PLACEHOLDER.TITLE: "title",
    PP_PLACEHOLDER.CENTER_TITLE: "title",
    PP_PLACEHOLDER.SUBTITLE: "subtitle",
    PP_PLACEHOLDER.BODY: "body",
    PP_PLACEHOLDER.OBJECT: "body",
    PP_PLACEHOLDER.TABLE: "table",
    PP_PLACEHOLDER.CHART: "chart",
    PP_PLACEHOLDER.BITMAP: "picture",
    PP_PLACEHOLDER.MEDIA_CLIP: "media",
    PP_PLACEHOLDER.ORG_CHART: "chart",
}


def analyze_master(master_pptx_path: str) -> dict:
    """Load a master .pptx template and extract its layout capabilities.

    Args:
        master_pptx_path: Absolute or relative path to the master PPTX file.

    Returns:
        A dict with key ``layouts``, each entry describing the layout name,
        index, and its placeholders (idx, type, dimensions).

    Raises:
        FileNotFoundError: If the template file does not exist.
        ValueError: If the file cannot be parsed as a valid PPTX.
    """
    path = Path(master_pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"Master template not found: {master_pptx_path}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ValueError(f"Cannot parse PPTX template: {exc}") from exc

    layouts: list[dict[str, Any]] = []

    for idx, layout in enumerate(prs.slide_layouts):
        placeholders: list[dict[str, Any]] = []

        for ph in layout.placeholders:
            ph_type = _resolve_placeholder_type(ph)
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": ph_type,
                "name": ph.name,
                "width_emu": ph.width,
                "height_emu": ph.height,
                "left_emu": ph.left,
                "top_emu": ph.top,
            })

        layouts.append({
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders,
        })

    return {"layouts": layouts, "slide_count": len(prs.slides)}


def _resolve_placeholder_type(ph) -> str:
    """Map a placeholder object to a human-readable type string."""
    ph_idx_type = ph.placeholder_format.type
    if ph_idx_type in _PLACEHOLDER_TYPE_MAP:
        return _PLACEHOLDER_TYPE_MAP[ph_idx_type]

    name_lower = ph.name.lower()
    if "picture" in name_lower or "image" in name_lower:
        return "picture"
    if "title" in name_lower:
        return "title"
    if "subtitle" in name_lower:
        return "subtitle"
    if "table" in name_lower:
        return "table"
    if "body" in name_lower or "content" in name_lower or "text" in name_lower:
        return "body"

    return "other"


def extract_content(pdf_path: str) -> dict:
    """Extract text and images from a PDF with correct reading order.

    Produces clean per-page text (column-aware) plus document-level font stats.
    The semantic classification (theory vs exercise) is delegated to the LLM,
    which is far more robust than font heuristics for educational PDFs.

    Args:
        pdf_path: Path to the source PDF.

    Returns:
        A dict with:
          - ``pages``: list of dicts, each with ``page``, ``text`` (reading-ordered)
          - ``full_text``: concatenated clean text across pages
          - ``images``: list of dicts with ``path``, ``page``, ``index``, dimensions
          - ``tmp_dir``: temp dir holding extracted images

    Raises:
        FileNotFoundError: If the PDF does not exist.
        RuntimeError: If PyMuPDF cannot open the file.
    """
    path = Path(pdf_path)
    if not path.exists():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    try:
        doc = fitz.open(str(path))
    except Exception as exc:
        raise RuntimeError(f"Cannot open PDF: {exc}") from exc

    pages: list[dict[str, Any]] = []
    images: list[dict[str, Any]] = []
    full_text_parts: list[str] = []
    tmp_dir = tempfile.mkdtemp(prefix="pptx_extract_")

    for page_num in range(len(doc)):
        page = doc[page_num]
        page_width = page.rect.width

        raw = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        text_blocks = [b for b in raw["blocks"] if b["type"] == 0]

        ordered_text = _reading_order_text(text_blocks, page_width)
        if ordered_text.strip():
            pages.append({"page": page_num + 1, "text": ordered_text})
            full_text_parts.append(ordered_text)

        for block in raw["blocks"]:
            if block["type"] == 1:
                w = block.get("width", 0) or (block["bbox"][2] - block["bbox"][0])
                h = block.get("height", 0) or (block["bbox"][3] - block["bbox"][1])
                if w < 60 or h < 60:
                    continue
                img_ext = block.get("ext", "png")
                img_filename = f"img_p{page_num + 1}_{uuid.uuid4().hex[:8]}.{img_ext}"
                img_path = os.path.join(tmp_dir, img_filename)
                try:
                    with open(img_path, "wb") as f:
                        f.write(block["image"])
                    images.append({
                        "path": img_path,
                        "page": page_num + 1,
                        "index": len(images),
                        "width": w,
                        "height": h,
                    })
                except (KeyError, IOError):
                    pass

    doc.close()

    return {
        "pages": pages,
        "full_text": "\n\n".join(full_text_parts),
        "images": images,
        "tmp_dir": tmp_dir,
    }


def _reading_order_text(text_blocks: list[dict], page_width: float) -> str:
    """Reconstruct page text in human reading order, handling 2-column layouts."""
    if not text_blocks:
        return ""

    mid_x = page_width / 2.0
    left_col = []
    right_col = []
    full_width = []

    for block in text_blocks:
        x0, _, x1, _ = block["bbox"]
        block_width = x1 - x0
        if block_width > page_width * 0.55:
            full_width.append(block)
        elif x1 <= mid_x + page_width * 0.05:
            left_col.append(block)
        elif x0 >= mid_x - page_width * 0.05:
            right_col.append(block)
        else:
            full_width.append(block)

    is_two_column = len(left_col) >= 2 and len(right_col) >= 2

    if is_two_column:
        ordered = (
            sorted(full_width, key=lambda b: b["bbox"][1])
            + sorted(left_col, key=lambda b: b["bbox"][1])
            + sorted(right_col, key=lambda b: b["bbox"][1])
        )
    else:
        ordered = sorted(text_blocks, key=lambda b: (b["bbox"][1], b["bbox"][0]))

    parts = [_extract_block_text(b) for b in ordered]
    return "\n".join(p for p in parts if p.strip())


def _extract_block_text(block: dict) -> str:
    """Concatenate span texts from a block, joining spans on the same line with spaces."""
    lines_out: list[str] = []
    for line in block.get("lines", []):
        spans = [s.get("text", "") for s in line.get("spans", [])]
        line_text = "".join(spans)
        if line_text.strip():
            lines_out.append(line_text.rstrip())
    return "\n".join(lines_out)


def _load_reference_data() -> tuple[dict, dict]:
    """Load the layout rules and few-shot examples derived from real operator decks."""
    base = Path(__file__).parent.parent
    rules, fewshot = {}, {}
    try:
        with open(base / "layout_rules.json", encoding="utf-8") as f:
            rules = json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        pass
    try:
        with open(base / "fewshot_examples.json", encoding="utf-8") as f:
            fewshot = json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        pass
    return rules, fewshot


def _build_layout_guide(master_layouts: dict | None, rules: dict) -> str:
    """Build a concise, operator-derived guide of which layout to use for what."""
    used_layouts = rules.get("layouts", {})
    available = {}
    if master_layouts:
        for layout in master_layouts.get("layouts", []):
            available[layout["name"]] = layout

    lines = []
    for name, info in used_layouts.items():
        match = available.get(name)
        idx_note = ""
        if match:
            idx_note = f" (layout_index={match['index']})"
        roles = []
        if info.get("title_placeholder_idx"):
            roles.append(f"title/heading idx={info['title_placeholder_idx']}")
        if info.get("heading_placeholder_idx"):
            roles.append(f"heading idx={info['heading_placeholder_idx']}")
        if info.get("body_placeholder_idx"):
            roles.append(f"body idx={info['body_placeholder_idx']}")
        if info.get("image_placeholder_idx"):
            roles.append(f"image idx={info['image_placeholder_idx']}")
        if info.get("table_placeholder_idx"):
            roles.append(f"table idx={info['table_placeholder_idx']}")
        lines.append(f'- "{name}"{idx_note}: {"; ".join(roles)}')
    return "\n".join(lines)


def _build_fewshot_block(fewshot: dict, max_examples: int = 2) -> str:
    """Build a compact few-shot block showing real PDF->slide deck structure."""
    examples = fewshot.get("examples", [])[:max_examples]
    blocks = []
    for ex in examples:
        slides_compact = []
        for s in ex.get("slides", []):
            ph = s.get("placeholders", {})
            slides_compact.append({
                "layout_name": s.get("layout_name"),
                "placeholders": ph,
            })
        blocks.append(f"UNIT {ex.get('unit')} (desired output deck):\n" +
                       json.dumps(slides_compact, ensure_ascii=False, indent=1))
    return "\n\n".join(blocks)


async def generate_variants(
    content: dict,
    master_layouts: dict | None,
    num_variants: int = 2,
    custom_prompt: str | None = None,
) -> list[dict[str, Any]]:
    """Generate a slide deck plan from PDF content, mimicking the expert operator.

    Produces, per content section/slide, ``num_variants`` layout alternatives.
    Follows the real operator grammar learned from training examples:
    keep only theory (discard exercises/audio/tests), rewrite concisely,
    use the 7 known layouts with the correct deck structure.

    Returns:
        List of section objects, each with a ``variants`` list.
    """
    if not USE_BEDROCK and not ANTHROPIC_API_KEY:
        raise EnvironmentError(
            "No LLM credentials configured. Set BEDROCK_AWS_ACCESS_KEY_ID + "
            "BEDROCK_AWS_SECRET_ACCESS_KEY (preferred) or ANTHROPIC_API_KEY."
        )

    rules, fewshot = _load_reference_data()
    layout_guide = _build_layout_guide(master_layouts, rules)
    fewshot_block = _build_fewshot_block(fewshot)

    pages = content.get("pages", [])
    if not pages and content.get("sections"):
        pages = [{"page": i + 1, "text": s.get("heading", "") + "\n" + s.get("text", "")}
                 for i, s in enumerate(content["sections"])]

    source_text = "\n\n".join(
        f"--- PDF page {p['page']} ---\n{p['text'][:3500]}" for p in pages
    )
    source_text = source_text[:45000]

    layout_index_map = {}
    if master_layouts:
        for layout in master_layouts.get("layouts", []):
            layout_index_map[layout["name"]] = layout["index"]

    user_message = f"""You are an expert editorial assistant for an Italian educational publisher. You convert English-language textbook unit PDFs into study/revision slide decks, following EXACTLY the style of the publisher's expert operator.

## HOW THE OPERATOR WORKS (learned from real examples)

1. SELECTION - This is critical. The PDF is a student workbook full of exercises, audio/vlog prompts, fill-in-the-gap tasks, true/false grids, "CRITICAL THINKING", "IN PAIRS", self-tests, vocabulary-match activities. The operator DISCARDS all of these. He keeps ONLY the THEORY: the declarative subject-matter content (definitions, classifications, processes, concepts) found in the reading boxes and numbered sub-sections.

2. REWRITE - The kept theory is REWRITTEN concisely: remove subordinate clauses, examples, hedging. Turn long passages into short declarative sentences. Lists become one short line per item. Target 8-40 words per slide body. NEVER copy the exercise instructions.

3. STRUCTURE (deck grammar) - The deck ALWAYS follows this skeleton:
   - Slide 1: "Title Slide" -> idx 0 = "Unit N\\n<Theme>" (theme from the unit opener page)
   - Each STEP in the PDF opens with "Capitolo + Immagine" -> idx 14 = "Step N\\n<Step name>"
   - Between dividers: content slides walking that step's theory sub-sections, in order
   - A single theory sub-section title (e.g. "Market research") becomes the HEADING, repeated across consecutive slides if the content is long, ALTERNATING layouts so the image flips side: 2_Testo -> 1_Testo -> Testo...

## LAYOUTS TO USE (use ONLY these, by exact name; pick layout_index from the map)
{layout_guide}

Layout name -> index map for THIS master: {json.dumps(layout_index_map, ensure_ascii=False)}

## LAYOUT SELECTION RULES
- Unit title -> "Title Slide" (idx 0)
- STEP divider -> "Capitolo + Immagine" (idx 14 = "Step N\\n<name>")
- A concept/definition with an image -> "2_Testo + 1 Immagine" (idx 15 heading, idx 18 body)
- Continuation of same concept (image flips) -> "1_Testo + 1 Immagine" (idx 15 body) then "Testo + 1 Immagine" (idx 15 heading, idx 18 body, no image)
- Two-sided comparison (X vs Y, Dos/Don'ts, hard/soft skills) -> "Tabella" (idx 15 heading, idx 14 table)
- Mind map / 3 parallel concepts -> "3 colonne Txt + Img" (idx 20 top heading, idx 15/17/19 columns)

## REAL EXAMPLES OF DESIRED OUTPUT
{fewshot_block}

## YOUR TASK
Read the source PDF text below. Produce the slide deck the operator would make.
For EACH slide, propose {num_variants} layout VARIANTS (alternative ways to lay out that same content - e.g. variant 1 uses "2_Testo + 1 Immagine", variant 2 uses "Testo + 1 Immagine"). The operator will pick one per slide.

Each variant's placeholder_fills must use the EXACT placeholder idx for that layout (from the rules above). For text use {{"type":"text","content":"..."}}. For images use {{"type":"image","suggestion":"<english search phrase>"}}. For tables use {{"type":"table","rows":[["col1","col2"],...]}}.

## SOURCE PDF
{source_text}

## OUTPUT (valid JSON only, no markdown fencing)
[
  {{
    "section_index": 0,
    "heading": "short label for this slide (for the operator UI)",
    "variants": [
      {{
        "layout_index": <int>,
        "layout_name": "<exact name>",
        "placeholder_fills": {{ "15": {{"type":"text","content":"..."}}, "18": {{"type":"text","content":"..."}}, "16": {{"type":"image","suggestion":"..."}} }},
        "design_rationale": "why this layout"
      }}
    ]
  }}
]
"""

    if custom_prompt:
        user_message += f"\n\n## ADDITIONAL OPERATOR INSTRUCTIONS\n{custom_prompt}"

    assistant_text = await _call_llm(user_message, max_tokens=32000)

    try:
        variants = json.loads(assistant_text)
    except json.JSONDecodeError as exc:
        cleaned = _attempt_json_extraction(assistant_text)
        if cleaned is None:
            raise RuntimeError(
                f"Claude returned invalid JSON. First 500 chars: {assistant_text[:500]}"
            ) from exc
        variants = cleaned

    return variants


async def _call_llm(prompt: str, max_tokens: int = 8192) -> str:
    """Route LLM call to Bedrock (preferred) or Anthropic API (fallback)."""
    if USE_BEDROCK:
        return await _call_bedrock(prompt, max_tokens)
    return await _call_anthropic_api(prompt, max_tokens)


async def _call_bedrock(prompt: str, max_tokens: int) -> str:
    """Call Claude via Amazon Bedrock using boto3."""
    import boto3
    from botocore.config import Config as BotoConfig

    # Re-read env at call time to support hot-reload
    region = os.environ.get("BEDROCK_AWS_REGION", "eu-west-1")
    access_key = os.environ.get("BEDROCK_AWS_ACCESS_KEY_ID", "")
    secret_key = os.environ.get("BEDROCK_AWS_SECRET_ACCESS_KEY", "")
    prefix = os.environ.get("BEDROCK_INFERENCE_PREFIX", "eu")
    model_id = os.environ.get("BEDROCK_MODEL", f"{prefix}.anthropic.claude-sonnet-4-6")

    boto_config = BotoConfig(
        read_timeout=300,
        connect_timeout=10,
        retries={"max_attempts": 1},
    )

    client = boto3.client(
        "bedrock-runtime",
        region_name=region,
        aws_access_key_id=access_key,
        aws_secret_access_key=secret_key,
        config=boto_config,
    )

    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": max_tokens,
        "temperature": 0.1,
        "messages": [{"role": "user", "content": [{"type": "text", "text": prompt}]}],
    })

    try:
        response = client.invoke_model(
            modelId=model_id,
            contentType="application/json",
            accept="application/json",
            body=body,
        )
    except Exception as exc:
        raise RuntimeError(f"Bedrock API call failed: {exc}") from exc

    response_body = json.loads(response["body"].read())
    return response_body["content"][0]["text"]


async def _call_anthropic_api(prompt: str, max_tokens: int) -> str:
    """Call Claude via direct Anthropic API (fallback if Bedrock not configured)."""
    payload = {
        "model": ANTHROPIC_MODEL,
        "max_tokens": max_tokens,
        "messages": [{"role": "user", "content": prompt}],
    }
    headers = {
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": ANTHROPIC_VERSION,
        "content-type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post(ANTHROPIC_API_URL, json=payload, headers=headers)
            response.raise_for_status()
    except httpx.HTTPStatusError as exc:
        raise RuntimeError(
            f"Anthropic API returned {exc.response.status_code}: "
            f"{exc.response.text[:500]}"
        ) from exc
    except httpx.RequestError as exc:
        raise RuntimeError(f"Anthropic API request failed: {exc}") from exc

    response_data = response.json()
    return response_data["content"][0]["text"]


def _default_layouts() -> list[dict]:
    """Fallback layout descriptions when no master template is provided."""
    return [
        {"index": 0, "name": "Title Slide", "placeholders": [
            {"idx": 0, "type": "title"}, {"idx": 1, "type": "subtitle"}
        ]},
        {"index": 1, "name": "Title and Content", "placeholders": [
            {"idx": 0, "type": "title"}, {"idx": 1, "type": "body"}
        ]},
        {"index": 2, "name": "Two Content", "placeholders": [
            {"idx": 0, "type": "title"}, {"idx": 1, "type": "body"}, {"idx": 2, "type": "body"}
        ]},
        {"index": 3, "name": "Title and Picture", "placeholders": [
            {"idx": 0, "type": "title"}, {"idx": 1, "type": "body"}, {"idx": 10, "type": "picture"}
        ]},
        {"index": 4, "name": "Blank with Image", "placeholders": [
            {"idx": 10, "type": "picture"}, {"idx": 1, "type": "body"}
        ]},
    ]


def _attempt_json_extraction(text: str):
    """Try to extract JSON from text that may contain markdown fencing."""
    for start_marker in ("```json", "```"):
        if start_marker in text:
            start = text.index(start_marker) + len(start_marker)
            end = text.find("```", start)
            if end == -1:
                end = len(text)
            try:
                return json.loads(text[start:end].strip())
            except (json.JSONDecodeError, ValueError):
                pass

    for i, ch in enumerate(text):
        if ch == "[":
            try:
                return json.loads(text[i:])
            except json.JSONDecodeError:
                pass
            break

    for i, ch in enumerate(text):
        if ch == "{":
            try:
                return [json.loads(text[i:])]
            except json.JSONDecodeError:
                pass
            break

    return None


def build_pptx(
    master_path: str,
    selected_variants: list[dict[str, Any]],
    images: dict[str, str] | None = None,
    output_path: str | None = None,
) -> str:
    """Assemble the final PPTX from a master template and selected variants.

    Args:
        master_path: Path to the master .pptx template.
        selected_variants: List of variant dicts (one per section), each specifying
            ``layout_index`` and ``placeholder_fills``.
        images: Optional mapping of placeholder idx (as string) to local image path.
            Overrides image suggestions with actual files.
        output_path: Where to save the result. If None, saves to a temp file.

    Returns:
        Absolute path to the generated PPTX file.

    Raises:
        FileNotFoundError: If master template or referenced images don't exist.
        ValueError: If a layout_index references a non-existent layout.
    """
    master = Path(master_path)
    if not master.exists():
        raise FileNotFoundError(f"Master template not found: {master_path}")

    prs = Presentation(str(master))
    images = images or {}

    for variant in selected_variants:
        layout_index = variant.get("layout_index", 0)
        if layout_index >= len(prs.slide_layouts):
            raise ValueError(
                f"Layout index {layout_index} out of range "
                f"(template has {len(prs.slide_layouts)} layouts)"
            )

        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

        fills = variant.get("placeholder_fills", {})
        for idx_str, fill_data in fills.items():
            idx = int(idx_str)
            _apply_fill(slide, idx, fill_data, images)

    if output_path is None:
        output_path = os.path.join(
            tempfile.gettempdir(),
            f"presentation_{uuid.uuid4().hex[:12]}.pptx",
        )

    prs.save(output_path)
    return os.path.abspath(output_path)


def _apply_fill(
    slide,
    placeholder_idx: int,
    fill_data: dict[str, Any],
    images: dict[str, str],
) -> None:
    """Apply a single fill instruction to a slide placeholder.

    Validates the placeholder exists; handles multi-paragraph text, images
    (left empty if no file, matching the operator workflow), and tables.
    """
    try:
        ph = slide.placeholders[placeholder_idx]
    except KeyError:
        return

    fill_type = fill_data.get("type", "text")
    content = fill_data.get("content", "")

    if fill_type == "text":
        _set_multiparagraph_text(ph, content)

    elif fill_type == "image":
        img_key = str(placeholder_idx)
        img_path = images.get(img_key) or fill_data.get("local_path")

        if img_path and Path(img_path).exists():
            try:
                ph.insert_picture(open(img_path, "rb"))
            except (AttributeError, Exception):
                try:
                    slide.shapes.add_picture(img_path, ph.left, ph.top, ph.width, ph.height)
                except Exception:
                    pass
        # If no image file: leave the placeholder empty (operator fills it later).
        # Never write "[Image: ...]" text into a picture slot.

    elif fill_type == "table":
        rows_data = fill_data.get("rows", [])
        if rows_data:
            _insert_table(slide, ph, rows_data)
        elif content:
            _set_multiparagraph_text(ph, content)


def _set_multiparagraph_text(ph, content: str) -> None:
    """Set text on a placeholder, preserving line breaks as separate paragraphs."""
    if not ph.has_text_frame:
        try:
            ph.text = content
        except Exception:
            pass
        return

    tf = ph.text_frame
    lines = content.split("\n") if content else [""]

    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line

    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _insert_table(slide, placeholder, rows: list[list[str]]) -> None:
    """Insert a table into the slide at the placeholder's position."""
    if not rows:
        return

    num_rows = len(rows)
    num_cols = max(len(row) for row in rows)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        placeholder.left, placeholder.top,
        placeholder.width, placeholder.height,
    )
    table = table_shape.table

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            if c_idx < num_cols:
                table.cell(r_idx, c_idx).text = str(cell_text)


# --- Approach C: Generate all slides as images for visual selection ---

LIBREOFFICE_PATH = os.environ.get(
    "LIBREOFFICE_PATH",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice"
)


def build_all_variants_pptx(
    master_path: str,
    variants_data: list[dict[str, Any]],
) -> str:
    """Build a single PPTX containing ALL variant slides for visual preview.

    For each section, adds N slides (one per variant), in order.
    Returns path to the temporary PPTX.
    """
    master = Path(master_path)
    if not master.exists():
        raise FileNotFoundError(f"Master template not found: {master_path}")

    prs = Presentation(str(master))

    # Remove any existing slides from the master template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    for section in variants_data:
        for variant in section.get("variants", []):
            layout_index = variant.get("layout_index", 0)
            if layout_index >= len(prs.slide_layouts):
                layout_index = 0

            layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(layout)

            fills = variant.get("placeholder_fills", {})
            for idx_str, fill_data in fills.items():
                _apply_fill(slide, int(idx_str), fill_data, {})

    output_path = os.path.join(
        tempfile.gettempdir(),
        f"all_variants_{uuid.uuid4().hex[:8]}.pptx",
    )
    prs.save(output_path)
    return output_path


def pptx_to_pngs(pptx_path: str, output_dir: str, dpi: int = 150) -> list[str]:
    """Convert a PPTX to individual PNG images per slide.

    Uses LibreOffice headless for PPTX→PDF, then PyMuPDF for PDF→PNG.
    Returns list of PNG file paths in slide order.
    """
    os.makedirs(output_dir, exist_ok=True)

    pdf_path = _pptx_to_pdf(pptx_path, output_dir)

    doc = fitz.open(pdf_path)
    png_paths = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)

        png_filename = f"slide_{page_num:03d}.png"
        png_path = os.path.join(output_dir, png_filename)
        pix.save(png_path)
        png_paths.append(png_path)

    doc.close()
    os.remove(pdf_path)
    return png_paths


def _pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """Convert PPTX to PDF using LibreOffice headless."""
    import subprocess

    soffice = LIBREOFFICE_PATH
    if not Path(soffice).exists():
        raise RuntimeError(
            f"LibreOffice not found at {soffice}. "
            "Install LibreOffice or set LIBREOFFICE_PATH."
        )

    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            pptx_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed: {result.stderr[:500]}"
        )

    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_name)

    if not Path(pdf_path).exists():
        existing = [f for f in os.listdir(output_dir) if f.endswith('.pdf')]
        if existing:
            pdf_path = os.path.join(output_dir, existing[0])
        else:
            raise RuntimeError("LibreOffice did not produce a PDF file")

    return pdf_path


async def generate_variants_with_thumbnails(
    content: dict,
    master_path: str,
    master_layouts: dict | None,
    num_variants: int = 2,
    custom_prompt: str | None = None,
    output_dir: str | None = None,
) -> dict[str, Any]:
    """Full pipeline: generate variants via Claude, build all slides, render to PNG.

    Returns a dict with:
    - generation_pptx_path: path to PPTX with all variant slides
    - sections: list of section dicts with thumbnail paths
    """
    variants_data = await generate_variants(
        content=content,
        master_layouts=master_layouts,
        num_variants=num_variants,
        custom_prompt=custom_prompt,
    )

    all_pptx_path = build_all_variants_pptx(master_path, variants_data)

    if output_dir is None:
        output_dir = os.path.join(tempfile.gettempdir(), f"thumbs_{uuid.uuid4().hex[:8]}")

    png_paths = pptx_to_pngs(all_pptx_path, output_dir)

    sections_result = []
    slide_idx = 0

    for section in variants_data:
        section_variants = []
        for v_idx, variant in enumerate(section.get("variants", [])):
            if slide_idx < len(png_paths):
                section_variants.append({
                    "variant_index": v_idx,
                    "slide_index": slide_idx,
                    "layout_index": variant.get("layout_index", 0),
                    "layout_name": variant.get("layout_name", ""),
                    "design_rationale": variant.get("design_rationale", ""),
                    "thumbnail_path": png_paths[slide_idx],
                })
                slide_idx += 1

        sections_result.append({
            "section_index": section.get("section_index", 0),
            "heading": section.get("heading", ""),
            "variants": section_variants,
        })

    return {
        "generation_pptx_path": all_pptx_path,
        "sections": sections_result,
    }


def build_final_from_selections(
    all_variants_pptx_path: str,
    selections: dict[int, int],
    sections_data: list[dict],
    output_path: str | None = None,
) -> str:
    """Build final PPTX by extracting only selected slides from the all-variants PPTX.

    Args:
        all_variants_pptx_path: Path to PPTX with all variant slides.
        selections: Mapping section_index -> variant_index (which variant was chosen).
        sections_data: The sections list from generate_variants_with_thumbnails.
        output_path: Where to save. If None, uses a temp file.

    Returns:
        Path to the final PPTX.
    """
    from copy import deepcopy
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    prs = Presentation(all_variants_pptx_path)

    selected_slide_indices = set()
    for section in sections_data:
        s_idx = section["section_index"]
        chosen_v = selections.get(s_idx, 0)
        for variant in section["variants"]:
            if variant["variant_index"] == chosen_v:
                selected_slide_indices.add(variant["slide_index"])
                break

    slides_to_remove = []
    for i in range(len(prs.slides)):
        if i not in selected_slide_indices:
            slides_to_remove.append(i)

    for idx in sorted(slides_to_remove, reverse=True):
        rId = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]

    if output_path is None:
        output_path = os.path.join(
            tempfile.gettempdir(),
            f"final_{uuid.uuid4().hex[:8]}.pptx",
        )

    prs.save(output_path)
    return output_path
